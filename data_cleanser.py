import cv2
import numpy as np
from PIL import Image
import os

# --- Presidio Imports for PII Redaction ---
from presidio_analyzer import AnalyzerEngine, PatternRecognizer
from presidio_anonymizer import AnonymizerEngine
from presidio_anonymizer.entities import OperatorConfig

# --- Text Extraction Logic (from our previous step) ---
import fitz  # PyMuPDF
import pytesseract
from pptx import Presentation

def extract_text_from_file(file_path):
    """Detects file type and extracts text. (Simplified from previous script)"""
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()
    
    if file_extension in ['.jpg', '.jpeg', '.png']:
        return pytesseract.image_to_string(Image.open(file_path))
    elif file_extension == '.pdf':
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text
    elif file_extension == '.pptx':
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return ""

# --- PART 1: TEXT CLEANSING WITH PRESIDIO ---

def anonymize_text_with_presidio(text_to_clean, client_name):
    """
    Analyzes text for PII and custom keywords (like client name) and redacts them.
    """
    # 1. Add a custom recognizer for the client's name
    client_name_recognizer = PatternRecognizer(
        supported_entity="CLIENT_NAME", 
        deny_list=[client_name]
    )
    
    # 2. Set up the analyzer
    analyzer = AnalyzerEngine()
    analyzer.registry.add_recognizer(client_name_recognizer)

    # 3. Analyze the text to find all PII and the client name
    analyzer_results = analyzer.analyze(text=text_to_clean, language='en')
    print(f"Found {len(analyzer_results)} sensitive items to redact.")

    # 4. Set up the anonymizer to replace entities
    anonymizer = AnonymizerEngine()
    
    # 5. Anonymize the text
    anonymized_result = anonymizer.anonymize(
        text=text_to_clean,
        analyzer_results=analyzer_results,
        operators={
            "DEFAULT": OperatorConfig("replace", {"new_value": "<REDACTED>"}),
            "CLIENT_NAME": OperatorConfig("replace", {"new_value": "<CLIENT_NAME_REDACTED>"})
        }
    )
    
    return anonymized_result.text

# --- PART 2: LOGO REDACTION WITH OPENCV ---

def redact_logo_from_image(main_image_path, logo_image_path, threshold=0.8):
    """
    Finds a logo in a main image and blacks it out.
    Returns an OpenCV image object.
    """
    main_img = cv2.imread(main_image_path)
    logo_img = cv2.imread(logo_image_path)
    
    if main_img is None or logo_img is None:
        print("Error: Could not read one of the images.")
        return None

    w, h = logo_img.shape[:-1]

    # Perform template matching
    res = cv2.matchTemplate(main_img, logo_img, cv2.TM_CCOEFF_NORMED)
    loc = np.where(res >= threshold)

    # Black out all found matches
    count = 0
    for pt in zip(*loc[::-1]):
        cv2.rectangle(main_img, pt, (pt[0] + h, pt[1] + w), (0, 0, 0), -1)
        count += 1
    
    print(f"Found and redacted {count} logo(s).")
    return main_img

# --- EXAMPLE USAGE ---
if __name__ == '__main__':
    # --- SETUP DUMMY FILES FOR TESTING ---
    if not os.path.exists('test_files'):
        os.makedirs('test_files')

    # Create a dummy text file with PII
    dummy_text = """
    This report is for ClientCorp.
    Please contact John Doe at 555-123-4567 for more details.
    The main office is at 123 Main St, Anytown, USA.
    """
    
    # Create a dummy main image and a dummy logo
    main_image = np.ones((400, 600, 3), dtype=np.uint8) * 255
    logo = np.zeros((50, 100, 3), dtype=np.uint8) # A black logo
    main_image[20:70, 50:150] = logo # Place logo on main image
    
    cv2.imwrite('test_files/document_with_logo.png', main_image)
    cv2.imwrite('test_files/client_logo.png', logo)

    # --- TEST THE CLEANSING PIPELINE ---

    # 1. Test Text Redaction
    print("--- Testing Text Redaction ---")
    client_to_redact = "ClientCorp"
    clean_text = anonymize_text_with_presidio(dummy_text, client_to_redact)
    print("\nOriginal Text:\n", dummy_text)
    print("\nCleaned Text:\n", clean_text)
    print("-" * 30)

    # 2. Test Logo Redaction
    print("\n--- Testing Logo Redaction ---")
    redacted_image_obj = redact_logo_from_image(
        'test_files/document_with_logo.png',
        'test_files/client_logo.png'
    )
    if redacted_image_obj is not None:
        cv2.imwrite('test_files/redacted_document.png', redacted_image_obj)
        print("Saved redacted image to 'test_files/redacted_document.png'")

    
