import streamlit as st
import os
import pandas as pd
from PIL import Image
import numpy as np
import cv2

# --- Import all our helper functions from the other files ---
# (For a real project, these would be in separate files as we planned.
# For this self-contained example, I'll include their core logic here.)

# --- 1. TEXT EXTRACTOR LOGIC ---
import fitz #PyMuPDF
import pytesseract
from pptx import Presentation

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_file(file_obj):
    file_extension = os.path.splitext(file_obj.name)[1].lower()
    if file_extension in ['.jpg', '.jpeg', '.png']:
        return pytesseract.image_to_string(Image.open(file_obj))
    elif file_extension == '.pdf':
        text = ""
        with fitz.open(stream=file_obj.read(), filetype="pdf") as doc:
            for page in doc:
                text += page.get_text()
        return text
    elif file_extension == '.pptx':
        text = ""
        prs = Presentation(file_obj)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return "Unsupported file format."

# --- 2. DATA CLEANSER LOGIC ---
from presidio_analyzer import AnalyzerEngine, PatternRecognizer
from presidio_anonymizer import AnonymizerEngine
from presidio_anonymizer.entities import OperatorConfig

def anonymize_text(text_to_clean, client_name):
    if not text_to_clean.strip() or not client_name.strip():
        return text_to_clean
    client_name_recognizer = PatternRecognizer(supported_entity="CLIENT_NAME", deny_list=[client_name])
    analyzer = AnalyzerEngine()
    analyzer.registry.add_recognizer(client_name_recognizer)
    analyzer_results = analyzer.analyze(text=text_to_clean, language='en')
    anonymizer = AnonymizerEngine()
    anonymized_result = anonymizer.anonymize(
        text=text_to_clean,
        analyzer_results=analyzer_results,
        operators={"DEFAULT": OperatorConfig("replace", {"new_value": "<REDACTED>"})}
    )
    return anonymized_result.text

def redact_logo(main_image, logo_image):
    main_img_cv = np.array(main_image.convert('RGB'))
    logo_img_cv = np.array(logo_image.convert('RGB'))
    w, h = logo_img_cv.shape[:-1]
    res = cv2.matchTemplate(main_img_cv, logo_img_cv, cv2.TM_CCOEFF_NORMED)
    threshold = 0.8
    loc = np.where(res >= threshold)
    for pt in zip(*loc[::-1]):
        cv2.rectangle(main_img_cv, pt, (pt[0] + h, pt[1] + w), (0, 0, 0), -1)
    return Image.fromarray(main_img_cv)

# --- 3. GROQ DATA EXTRACTOR & SUMMARIZER LOGIC ---
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.pydantic_v1 import BaseModel, Field
from typing import List

class KeyFinding(BaseModel):
    finding: str = Field(description="A single, crucial key finding from the text.")

class FileSummary(BaseModel):
    file_description: str = Field(description="A brief, neutral, one or two-sentence description of the document's content.")
    key_findings: List[KeyFinding] = Field(description="A list of the most important key findings, such as firewall rules, IAM policies, or security vulnerabilities.")

def get_summary_and_findings(text_to_analyze):
    if not os.getenv("GROQ_API_KEY"):
        st.error("GROQ_API_KEY environment variable not set!")
        return None
    if not text_to_analyze.strip():
        return {"file_description": "No text could be extracted from the file.", "key_findings": []}
        
    parser = JsonOutputParser(pydantic_object=FileSummary)
    prompt = ChatPromptTemplate.from_messages([
        ("system", "You are an expert security analyst. Extract a file description and key findings from the user's text. Respond with a JSON object that follows this format: {format_instructions}"),
        ("user", "{text}")
    ])
    llm = ChatGroq(temperature=0, model_name="llama-3.1-8b-instant")
    chain = prompt | llm | parser
    return chain.invoke({"text": text_to_analyze, "format_instructions": parser.get_format_instructions()})


# --- STREAMLIT UI ---
st.set_page_config(layout="wide")
st.title("🔒 Automated File Cleanser and Analyzer")

st.info("This tool automates the cleansing of sensitive data from files and extracts key security insights, presenting them in a final summary table.")

# --- INPUTS ---
col1, col2 = st.columns(2)
with col1:
    uploaded_file = st.file_uploader("1. Upload a document to analyze", type=['pdf', 'png', 'jpeg', 'jpg', 'pptx'])
    client_name = st.text_input("2. Enter Client Name to Redact", placeholder="e.g., ClientCorp")

with col2:
    logo_file = st.file_uploader("3. (Optional) Upload a client logo to redact from images", type=['png', 'jpeg', 'jpg'])


# --- PROCESSING AND OUTPUT ---
if st.button("Process File", type="primary") and uploaded_file:
    with st.spinner("Processing... Please wait."):
        
        # --- Step 1: Extract Text ---
        st.write("Step 1: Extracting text from file...")
        raw_text = extract_text_from_file(uploaded_file)
        
        # --- Step 2: Cleanse Data ---
        st.write("Step 2: Cleansing text and redacting PII...")
        anonymized_text = anonymize_text(raw_text, client_name)
        
        final_image = None
        if uploaded_file.type.startswith('image/') and logo_file:
            st.write("Step 2b: Redacting logo from image...")
            original_image = Image.open(uploaded_file)
            logo_image = Image.open(logo_file)
            final_image = redact_logo(original_image, logo_image)

        # --- Step 3: Analyze and Summarize ---
        st.write("Step 3: Generating description and key findings with AI...")
        analysis_result = get_summary_and_findings(anonymized_text)

        # --- Step 4: Display Final Table ---
        st.write("Step 4: Generating final report...")
        if analysis_result:
            file_name = uploaded_file.name
            file_type = os.path.splitext(file_name)[1]
            description = analysis_result['file_description']
            
            # Format key findings with bullet points
            findings_list = [f"- {item['finding']}" for item in analysis_result.get('key_findings', [])]
            key_findings_str = "\n".join(findings_list)
            
            # Create DataFrame for the table
            summary_data = {
                "File Name": [file_name],
                "File Type": [file_type],
                "File Description": [description],
                "Key Findings": [key_findings_str]
            }
            df = pd.DataFrame(summary_data)
            
            st.success("✅ Processing Complete!")
            st.subheader("Final Summary Report")
            st.dataframe(df, use_container_width=True)

            if final_image:
                st.subheader("Redacted Image")
                st.image(final_image, caption="Image with logo redacted.")
        else:
            st.error("Failed to generate analysis. Please check your API key and try again.")

else:
    st.warning("Please upload a file and click 'Process File' to begin.")
