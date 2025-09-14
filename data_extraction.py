import os
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pptx import Presentation

# --- HELPER FUNCTIONS FOR EACH FILE TYPE ---

def extract_text_from_image(file_path):
    """
    Uses Tesseract OCR to extract text from an image file (PNG, JPG).
    """
    try:
        return pytesseract.image_to_string(Image.open(file_path))
    except Exception as e:
        return f"Error processing image {file_path}: {e}"

def extract_text_from_pdf(file_path):
    """
    Extracts text from a PDF file. It handles both text-based and image-based (scanned) PDFs.
    """
    text = ""
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            # First, try to extract text directly
            page_text = page.get_text()
            if page_text:
                text += page_text
            else:
                # If no text, assume it's a scanned image and use OCR
                try:
                    pix = page.get_pixmap()
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text += pytesseract.image_to_string(img)
                except Exception as ocr_error:
                    print(f"Could not perform OCR on page {page_num} of {file_path}: {ocr_error}")
        doc.close()
        return text
    except Exception as e:
        return f"Error processing PDF {file_path}: {e}"

def extract_text_from_pptx(file_path):
    """
    Extracts text from all slides in a PowerPoint (PPTX) file.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error processing PPTX {file_path}: {e}"

# --- THE MAIN ROUTER FUNCTION ---

def extract_text_from_file(file_path):
    """
    Detects the file type and calls the appropriate helper function to extract text.
    """
    # Get the file extension and convert to lowercase
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()
    
    print(f"Processing file: {file_path} (Type: {file_extension})")

    if file_extension in ['.jpg', '.jpeg', '.png']:
        return extract_text_from_image(file_path)
    elif file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    else:
        return f"Unsupported file type: {file_extension}"

# --- EXAMPLE USAGE ---
if __name__ == '__main__':
    # Create dummy files for testing purposes
    # In your actual project, you would use the uploaded file's path
    if not os.path.exists('test_files'):
        os.makedirs('test_files')

    # Create a dummy image file (requires Pillow)
    try:
        img = Image.new('RGB', (400, 100), color = 'white')
        img.save('test_files/test_image.png')
        print("Created dummy image.")
    except Exception as e:
        print(f"Could not create dummy image: {e}")

    # Create a dummy pptx file (requires python-pptx)
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "This is a test slide."
        prs.save('test_files/test_presentation.pptx')
        print("Created dummy presentation.")
    except Exception as e:
        print(f"Could not create dummy presentation: {e}")


    # --- Test the extractor with different files ---
    # Replace these with the actual paths to your files
    files_to_test = [
        'test_files/test_image.png',
        'test_files/test_presentation.pptx',
        # 'path/to/your/sample.pdf', # Add a PDF to test
        'test_files/unsupported_file.txt'
    ]
    
    # Create a dummy txt file for the unsupported case
    with open('test_files/unsupported_file.txt', 'w') as f:
        f.write("This is a test.")

    for file in files_to_test:
        if os.path.exists(file):
            print("-" * 20)
            extracted_text = extract_text_from_file(file)
            print("--- Extracted Text ---")
            print(extracted_text[:300] + "..." if extracted_text else "No text found.") # Print first 300 chars
            print("-" * 20 + "\n")
        else:
            print(f"File not found: {file}")
