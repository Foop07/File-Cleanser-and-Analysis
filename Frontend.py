
import streamlit as st
import os
import pandas as pd
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO

# --- Import backend modules ---
from utils.text_extractor import extract_text
from utils.data_cleanser import anonymize_text
from utils.content_interpreter import identify_client_name, interpret_content

# --- UI Rendering Function ---
def display_report(file, result):
    st.header(f"Intelligence Report: {file.name}")

    col1, col2 = st.columns(2)
    col1.metric("File Type", os.path.splitext(file.name)[1])
    pii_level = result.get('pii_sensitivity_level', 'N/A')
    col2.metric("PII Sensitivity Level", pii_level)
    
    st.markdown("---")
    st.subheader("Executive Summary")
    st.write(result.get('executive_summary', 'No summary available.'))
    
    st.subheader("Detailed Findings")
    findings = result.get('detailed_findings', {})
    if not findings or all(not v for v in findings.values()):
        st.info("No detailed findings were extracted from this document.")
    else:
        for category, items in findings.items():
            if items:
                st.markdown(f"**{category.replace('_', ' ').title()}**")
                for item in items:
                    with st.expander(f"{item['title']}"):
                        st.write(item['description'])

# --- Main Application UI ---
st.set_page_config(layout="wide", page_title="Intelligent File Analyzer")
st.title("🚀 Intelligent File Analyzer")
st.info("Upload documents. The AI will identify the client, redact data, and generate detailed reports.")

uploaded_files = st.file_uploader(
    "Upload Documents for Analysis",
    type=['pdf', 'png', 'jpeg', 'jpg', 'pptx', 'xlsx', 'txt', 'csv'],
    accept_multiple_files=True
)

if st.button("Generate Intelligence Reports", type="primary"):
    if not uploaded_files:
        st.warning("Please upload at least one file to begin analysis.")
    else:
        st.header("Batch Analysis Results")
        
        for file in uploaded_files:
            st.markdown("---")
            with st.spinner(f"Processing {file.name}..."):
                raw_text = ""
                df_to_save = None
                pptx_redacted = None
                error = None
                ext = os.path.splitext(file.name)[1].lower()

                # --- TEXT EXTRACTION ---
                try:
                    if ext == ".txt":
                        raw_text = file.read().decode("utf-8")
                    elif ext == ".csv":
                        df = pd.read_csv(file)
                        raw_text = df.to_csv(index=False)
                        df_to_save = df.copy()
                    elif ext == ".xlsx":
                        df = pd.read_excel(file)
                        raw_text = df.to_csv(index=False)
                        df_to_save = df.copy()
                    elif ext == ".pdf":
                        doc = fitz.open(stream=file.read(), filetype="pdf")
                        raw_text = ""
                        for page in doc:
                            raw_text += page.get_text()
                    elif ext == ".pptx":
                        prs = Presentation(file)
                        raw_text = ""
                        pptx_redacted = prs
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    raw_text += shape.text + "\n"
                    else:
                        extraction_result = extract_text(file)
                        raw_text = extraction_result.get("text", "")
                        error = extraction_result.get("error")

                    if not raw_text.strip() and not error:
                        error = "No text could be extracted."
                except Exception as e:
                    error = str(e)

                if error:
                    st.error(f"**{file.name}**: Could not process file. Reason: {error}")
                    raw_text = f"Could not extract content from {file.name}. Error: {error}"

                # --- CLIENT IDENTIFICATION & ANONYMIZATION ---
                try:
                    client_name = identify_client_name(raw_text)
                except Exception as e:
                    client_name = "Unknown"
                    st.warning(f"Could not identify client: {e}")
                
                st.write(f"🤖 AI identified client for **{file.name}** as: **{client_name}**")
                
                try:
                    anonymized_text, pii_results = anonymize_text(raw_text, client_name)
                    pii_count = len(pii_results)
                except Exception as e:
                    anonymized_text = raw_text
                    pii_count = 0
                    st.warning(f"Anonymization failed: {e}")

                # --- CONTENT INTERPRETATION ---
                try:
                    analysis_result = interpret_content(anonymized_text, pii_count)
                    if "error" in analysis_result:
                        st.error(f"Could not analyze {file.name}: {analysis_result['error']}")
                except Exception as e:
                    analysis_result = {}
                    st.warning(f"Content interpretation failed: {e}")

                # --- DISPLAY REPORT ---
                display_report(file, analysis_result)

                # --- DOWNLOAD BUTTON ---
                if ext in [".xlsx", ".csv"] and df_to_save is not None:
                    # Redact all text cells in DataFrame
                    df_redacted = df_to_save.copy()
                    for col in df_redacted.columns:
                        df_redacted[col] = df_redacted[col].astype(str).apply(lambda x: anonymize_text(x, client_name)[0])

                    if ext == ".xlsx":
                        output_file_name = f"redacted_{os.path.splitext(file.name)[0]}.xlsx"
                        df_redacted.to_excel(output_file_name, index=False)
                        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    else:
                        output_file_name = f"redacted_{os.path.splitext(file.name)[0]}.csv"
                        df_redacted.to_csv(output_file_name, index=False)
                        mime_type = "text/csv"

                    with open(output_file_name, "rb") as f:
                        st.download_button(
                            label="Download Redacted File",
                            data=f,
                            file_name=output_file_name,
                            mime=mime_type
                        )

                elif ext == ".pptx" and pptx_redacted is not None:
                    # Replace text in slides with redacted version
                    prs = pptx_redacted
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                redacted_text, _ = anonymize_text(shape.text, client_name)
                                shape.text = redacted_text

                    output_file_name = f"redacted_{os.path.splitext(file.name)[0]}.pptx"
                    prs.save(output_file_name)

                    with open(output_file_name, "rb") as f:
                        st.download_button(
                            label="Download Redacted File",
                            data=f,
                            file_name=output_file_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                else:
                    # TXT, PDF, images
                    output_file_name = f"redacted_{os.path.splitext(file.name)[0]}.txt"
                    with open(output_file_name, "w", encoding="utf-8") as f:
                        f.write(anonymized_text)
                    with open(output_file_name, "rb") as f:
                        st.download_button(
                            label="Download Redacted File",
                            data=f,
                            file_name=output_file_name,
                            mime="text/plain"
                        )

        st.success("✅ Batch processing complete.")
