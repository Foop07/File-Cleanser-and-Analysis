import os
from PIL import Image
import pytesseract
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd
import cv2 # OpenCV for image processing
import numpy as np
import streamlit as st

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# --- NEW: Image Pre-processing Function ---
def preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    """Applies several pre-processing techniques to an image to improve OCR accuracy."""
    # Convert the Pillow image to an OpenCV-compatible format (NumPy array)
    img_cv = np.array(image.convert('RGB'))
    
    # 1. Convert to grayscale, as color information is not needed for OCR
    gray = cv2.cvtColor(img_cv, cv2.COLOR_RGB2GRAY)
    
    # 2. Apply a threshold to create a binary (black and white) image.
    # This is the most critical step for separating text from the background.
    _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)
    
    # 3. (Optional but helpful) Upscale the image if it's small. Tesseract works
    #    best on images with a DPI of at least 300, so larger text is better.
    h, w = thresh.shape
    if h < 500 or w < 500: # Heuristic for small images
        thresh = cv2.resize(thresh, (w*2, h*2), interpolation=cv2.INTER_CUBIC)
        
    # Convert the processed OpenCV image back to a Pillow image for pytesseract
    return Image.fromarray(thresh)


# --- Updated Extraction Function ---
def extract_text(uploaded_file: st.runtime.uploaded_file_manager.UploadedFile) -> dict:
    """
    Detects a file's type and extracts text, with advanced pre-processing for images.
    Returns a dictionary to handle potential errors gracefully.
    """
    try:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        
        # --- THIS IS THE UPGRADED PART ---
        if file_extension in ['.jpg', '.jpeg', '.png']:
            image = Image.open(uploaded_file)
            # Add the new pre-processing step before performing OCR
            processed_image = preprocess_image_for_ocr(image)
            text = pytesseract.image_to_string(processed_image)
            return {"text": text, "error": None}

        # --- (Rest of the file types remain the same) ---
        elif file_extension == '.pdf':
            text = ""
            with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
            return {"text": text, "error": None}

        elif file_extension == '.pptx':
            text = ""
            prs = Presentation(uploaded_file)
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Extract text from text frames
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                text += " | ".join(row_text) + "\n"
                    
                    # Extract text from grouped shapes
                    if hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text:
                                text += sub_shape.text + "\n"
                
                text += "\n"  # Add spacing between slides
            
            return {"text": text, "error": None}

        elif file_extension == '.xlsx':
            text = ""
            excel_sheets = pd.read_excel(uploaded_file, sheet_name=None)
            for sheet_name, df in excel_sheets.items():
                text += f"--- Sheet: {sheet_name} ---\n{df.to_string()}\n\n"
            return {"text": text, "error": None}
            
        else:
            return {"text": None, "error": f"Unsupported file format: {file_extension}"}

    except Exception as e:
        return {"text": None, "error": f"An error occurred while processing {uploaded_file.name}: {e}"}

